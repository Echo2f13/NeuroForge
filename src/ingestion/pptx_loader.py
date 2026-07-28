"""PPTX Loader for NeuroForge.

Extracts text content from PowerPoint (.pptx) files using python-pptx.
Handles slide text (titles, bullets, text boxes), speaker notes,
embedded image detection, and maintains slide ordering.

Usage:
    from src.ingestion.pptx_loader import PPTXLoader

    loader = PPTXLoader()
    doc = loader.load("path/to/presentation.pptx")

    # Legacy function interface still available:
    from src.ingestion.pptx_loader import extract_pptx
    doc = extract_pptx("path/to/presentation.pptx")
"""

from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.shapes.base import BaseShape
from pptx.slide import Slide

sys.path.insert(0, str(Path(__file__).resolve().parents[2]))

from models.document import Document, DocumentMetadata, InputFormat, Section


def extract_shape_text(shape: BaseShape) -> str:
    """Extract all text from a single slide shape.

    Handles text frames (titles, body placeholders, text boxes)
    by iterating through paragraphs and runs.

    Args:
        shape: A python-pptx shape object.

    Returns:
        The concatenated text content of the shape, with paragraphs
        separated by newlines. Returns empty string if shape has no text.
    """
    if not shape.has_text_frame:
        return ""

    paragraphs = []
    for paragraph in shape.text_frame.paragraphs:
        text = paragraph.text.strip()
        if text:
            paragraphs.append(text)

    return "\n".join(paragraphs)


def _extract_slide_content(slide: Slide, slide_number: int) -> tuple[str | None, str, list[str]]:
    """Extract title, body text, and image markers from a slide.

    Args:
        slide: A python-pptx Slide object.
        slide_number: 1-based slide number for image annotations.

    Returns:
        A tuple of (title, body_text, image_markers) where:
        - title: The slide title text or None if no title placeholder.
        - body_text: All non-title text content joined by newlines.
        - image_markers: List of image annotation strings.
    """
    title: str | None = None
    body_parts: list[str] = []
    image_markers: list[str] = []

    for shape in slide.shapes:
        # Detect embedded images
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            image_markers.append(f"[Image on slide {slide_number}]")
            continue

        # Check for group shapes that may contain images
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            # Groups can contain pictures; note them
            image_markers.append(f"[Image on slide {slide_number}]")

        # Extract title from title placeholder
        if shape.has_text_frame:
            is_title = False
            # Check if shape is a placeholder by inspecting the XML element
            if shape._element.ph is not None:
                # ph idx 0 is the title placeholder
                idx = shape._element.ph.get("idx", "0")
                if idx == "0" or idx == 0:
                    is_title = True

            if is_title:
                title_text = extract_shape_text(shape)
                if title_text:
                    title = title_text
            else:
                text = extract_shape_text(shape)
                if text:
                    body_parts.append(text)

    body_text = "\n".join(body_parts)
    return title, body_text, image_markers


def _extract_speaker_notes(slide: Slide) -> str:
    """Extract speaker notes from a slide.

    Args:
        slide: A python-pptx Slide object.

    Returns:
        The speaker notes text, or empty string if none present.
    """
    if slide.has_notes_slide:
        notes_slide = slide.notes_slide
        notes_tf = notes_slide.notes_text_frame
        if notes_tf and notes_tf.text.strip():
            return notes_tf.text.strip()
    return ""


def extract_pptx(file_path: str) -> Document:
    """Extract content from a PPTX file into a unified Document.

    Opens the specified PowerPoint file and extracts text from each slide
    including titles, body text, text boxes, and speaker notes. Maintains
    slide ordering and notes presence of embedded images in metadata.

    Args:
        file_path: Path to the .pptx file.

    Returns:
        A Document object containing:
        - Full concatenated text content
        - Metadata (source path, format=PPTX, slide count, title)
        - Sections (one per slide with heading=title, content=body+notes)

    Raises:
        FileNotFoundError: If the file does not exist.
        ValueError: If the file is not a valid .pptx file or is corrupt.
    """
    path = Path(file_path)

    if not path.exists():
        raise FileNotFoundError(f"PPTX file not found: {file_path}")

    if not path.suffix.lower() == ".pptx":
        raise ValueError(f"File is not a .pptx file: {file_path}")

    try:
        prs = Presentation(str(path))
    except Exception as e:
        raise ValueError(f"Failed to open PPTX file (corrupt or invalid): {file_path}. Error: {e}")

    sections: list[Section] = []
    all_content_parts: list[str] = []
    document_title: str | None = None

    for slide_number, slide in enumerate(prs.slides, start=1):
        title, body_text, image_markers = _extract_slide_content(slide, slide_number)
        notes = _extract_speaker_notes(slide)

        # Use first slide's title as document title
        if slide_number == 1 and title:
            document_title = title

        # Build section content
        section_parts: list[str] = []
        if body_text:
            section_parts.append(body_text)
        if image_markers:
            section_parts.extend(image_markers)
        if notes:
            section_parts.append(f"Speaker Notes: {notes}")

        section_content = "\n".join(section_parts) if section_parts else ""

        # Only create a section if there's meaningful content
        if title or section_content:
            # Section content must be non-empty per model constraint
            final_content = section_content if section_content else title or ""
            sections.append(
                Section(
                    heading=title,
                    content=final_content,
                    level=1,
                    page_number=slide_number,
                )
            )

        # Build full document content
        slide_text_parts: list[str] = []
        if title:
            slide_text_parts.append(title)
        if body_text:
            slide_text_parts.append(body_text)
        if image_markers:
            slide_text_parts.extend(image_markers)
        if notes:
            slide_text_parts.append(f"Speaker Notes: {notes}")

        if slide_text_parts:
            all_content_parts.append("\n".join(slide_text_parts))

    full_content = "\n\n".join(all_content_parts) if all_content_parts else ""

    # Ensure content is non-empty (model requires min_length=1)
    if not full_content:
        full_content = "[Empty presentation]"

    total_slides = len(prs.slides)

    metadata = DocumentMetadata(
        source=str(path.resolve()),
        format=InputFormat.PPTX,
        title=document_title,
        total_pages=total_slides if total_slides > 0 else None,
    )

    return Document(
        content=full_content,
        metadata=metadata,
        sections=sections,
    )


class PPTXLoader:
    """Loader class for extracting content from PowerPoint (.pptx) files.

    Provides a class-based interface for PPTX ingestion. Extracts slide text
    (titles, bullets, text boxes), speaker notes, and detects embedded images.
    Maintains slide ordering throughout.

    Usage:
        loader = PPTXLoader()
        doc = loader.load("path/to/presentation.pptx")
    """

    def load(self, file_path: str) -> Document:
        """Load and extract content from a PPTX file.

        Opens the specified PowerPoint file and extracts text from each slide
        including titles, body text, text boxes, and speaker notes. Maintains
        slide ordering and notes presence of embedded images in metadata.

        Each slide becomes a Section with:
        - heading = slide title (or None if no title)
        - content = body text + image markers + speaker notes
        - page_number = slide number (1-based)

        Args:
            file_path: Path to the .pptx file.

        Returns:
            A Document object containing:
            - Full concatenated text content
            - Metadata (source path, format=PPTX, slide count, title)
            - Sections (one per slide with heading=title, content=body+notes)

        Raises:
            FileNotFoundError: If the file does not exist.
            ValueError: If the file is not a valid .pptx file or is corrupt.
        """
        return extract_pptx(file_path)
