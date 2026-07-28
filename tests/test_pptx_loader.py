"""Tests for the PPTX loader module."""

import os
import sys
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches

# Ensure project root is importable
sys.path.insert(0, str(Path(__file__).resolve().parents[1]))

from models.document import InputFormat
from src.ingestion.pptx_loader import PPTXLoader, extract_pptx, extract_shape_text


@pytest.fixture
def sample_pptx(tmp_path: Path) -> Path:
    """Create a sample PPTX file for testing."""
    prs = Presentation()

    # Slide 1: Title slide with title and subtitle
    slide_layout = prs.slide_layouts[0]  # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Test Presentation"
    slide.placeholders[1].text = "A subtitle for testing"

    # Slide 2: Content slide with bullets and speaker notes
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Slide Two"
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "Bullet point 1"
    p = tf.add_paragraph()
    p.text = "Bullet point 2"
    p = tf.add_paragraph()
    p.text = "Bullet point 3"

    # Add speaker notes to slide 2
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = "These are speaker notes for slide 2."

    # Slide 3: Blank slide with a text box
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(2))
    tf = txBox.text_frame
    tf.text = "Custom text box content"

    file_path = tmp_path / "test_presentation.pptx"
    prs.save(str(file_path))
    return file_path


@pytest.fixture
def pptx_with_image(tmp_path: Path) -> Path:
    """Create a PPTX file containing an embedded image."""
    prs = Presentation()

    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Add title via text box
    txBox = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(5), Inches(1))
    txBox.text_frame.text = "Slide With Image"

    # Create a small test image
    img_path = tmp_path / "test_img.png"
    # Create a minimal valid PNG file (1x1 pixel)
    import struct
    import zlib

    def create_minimal_png():
        signature = b"\x89PNG\r\n\x1a\n"
        # IHDR chunk
        ihdr_data = struct.pack(">IIBBBBB", 1, 1, 8, 2, 0, 0, 0)
        ihdr_crc = zlib.crc32(b"IHDR" + ihdr_data) & 0xFFFFFFFF
        ihdr = struct.pack(">I", 13) + b"IHDR" + ihdr_data + struct.pack(">I", ihdr_crc)
        # IDAT chunk
        raw_data = zlib.compress(b"\x00\xff\x00\x00")
        idat_crc = zlib.crc32(b"IDAT" + raw_data) & 0xFFFFFFFF
        idat = struct.pack(">I", len(raw_data)) + b"IDAT" + raw_data + struct.pack(">I", idat_crc)
        # IEND chunk
        iend_crc = zlib.crc32(b"IEND") & 0xFFFFFFFF
        iend = struct.pack(">I", 0) + b"IEND" + struct.pack(">I", iend_crc)
        return signature + ihdr + idat + iend

    img_path.write_bytes(create_minimal_png())
    slide.shapes.add_picture(str(img_path), Inches(2), Inches(2), Inches(3), Inches(2))

    file_path = tmp_path / "image_presentation.pptx"
    prs.save(str(file_path))
    return file_path


@pytest.fixture
def empty_pptx(tmp_path: Path) -> Path:
    """Create an empty PPTX file (no slides)."""
    prs = Presentation()
    file_path = tmp_path / "empty_presentation.pptx"
    prs.save(str(file_path))
    return file_path


class TestExtractPptx:
    """Tests for the extract_pptx function."""

    def test_extracts_document_with_correct_metadata(self, sample_pptx: Path):
        doc = extract_pptx(str(sample_pptx))

        assert doc.metadata.format == InputFormat.PPTX
        assert doc.metadata.title == "Test Presentation"
        assert doc.metadata.total_pages == 3
        assert str(sample_pptx.resolve()) in doc.metadata.source

    def test_extracts_slide_titles(self, sample_pptx: Path):
        doc = extract_pptx(str(sample_pptx))

        headings = [s.heading for s in doc.sections]
        assert "Test Presentation" in headings
        assert "Slide Two" in headings

    def test_extracts_body_text(self, sample_pptx: Path):
        doc = extract_pptx(str(sample_pptx))

        # Find slide 2 section
        slide_two = next(s for s in doc.sections if s.heading == "Slide Two")
        assert "Bullet point 1" in slide_two.content
        assert "Bullet point 2" in slide_two.content
        assert "Bullet point 3" in slide_two.content

    def test_extracts_speaker_notes(self, sample_pptx: Path):
        doc = extract_pptx(str(sample_pptx))

        # Slide 2 has speaker notes
        slide_two = next(s for s in doc.sections if s.heading == "Slide Two")
        assert "Speaker Notes:" in slide_two.content
        assert "speaker notes for slide 2" in slide_two.content

    def test_extracts_text_boxes(self, sample_pptx: Path):
        doc = extract_pptx(str(sample_pptx))

        # Slide 3 has a text box
        assert "Custom text box content" in doc.content

    def test_maintains_slide_ordering(self, sample_pptx: Path):
        doc = extract_pptx(str(sample_pptx))

        # Sections should be ordered by page_number
        page_numbers = [s.page_number for s in doc.sections]
        assert page_numbers == sorted(page_numbers)

    def test_detects_embedded_images(self, pptx_with_image: Path):
        doc = extract_pptx(str(pptx_with_image))

        assert "[Image on slide 1]" in doc.content

    def test_handles_empty_presentation(self, empty_pptx: Path):
        doc = extract_pptx(str(empty_pptx))

        assert doc.content == "[Empty presentation]"
        assert doc.metadata.total_pages is None or doc.metadata.total_pages >= 1
        assert doc.sections == []

    def test_file_not_found_raises_error(self):
        with pytest.raises(FileNotFoundError):
            extract_pptx("nonexistent_file.pptx")

    def test_invalid_extension_raises_error(self, tmp_path: Path):
        txt_file = tmp_path / "not_a_pptx.txt"
        txt_file.write_text("This is not a PPTX file")

        with pytest.raises(ValueError, match="not a .pptx file"):
            extract_pptx(str(txt_file))

    def test_corrupt_file_raises_error(self, tmp_path: Path):
        corrupt_file = tmp_path / "corrupt.pptx"
        corrupt_file.write_bytes(b"not a valid pptx content")

        with pytest.raises(ValueError, match="corrupt or invalid"):
            extract_pptx(str(corrupt_file))

    def test_sections_have_page_numbers(self, sample_pptx: Path):
        doc = extract_pptx(str(sample_pptx))

        for section in doc.sections:
            assert section.page_number is not None
            assert section.page_number >= 1


class TestExtractShapeText:
    """Tests for the extract_shape_text helper function."""

    def test_extracts_text_from_shape(self, sample_pptx: Path):
        prs = Presentation(str(sample_pptx))
        slide = prs.slides[0]
        # Title shape
        title_shape = slide.shapes.title
        text = extract_shape_text(title_shape)
        assert text == "Test Presentation"

    def test_returns_empty_for_no_text_frame(self, pptx_with_image: Path):
        prs = Presentation(str(pptx_with_image))
        slide = prs.slides[0]
        # Find the picture shape
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                text = extract_shape_text(shape)
                assert text == ""
                break


class TestPPTXLoader:
    """Tests for the PPTXLoader class interface."""

    def test_load_returns_document(self, sample_pptx: Path):
        loader = PPTXLoader()
        doc = loader.load(str(sample_pptx))

        assert doc.metadata.format == InputFormat.PPTX
        assert doc.metadata.title == "Test Presentation"
        assert doc.metadata.total_pages == 3

    def test_load_extracts_all_slides(self, sample_pptx: Path):
        loader = PPTXLoader()
        doc = loader.load(str(sample_pptx))

        # Should have sections for slides with content
        assert len(doc.sections) >= 2
        headings = [s.heading for s in doc.sections]
        assert "Test Presentation" in headings
        assert "Slide Two" in headings

    def test_load_maintains_slide_ordering(self, sample_pptx: Path):
        loader = PPTXLoader()
        doc = loader.load(str(sample_pptx))

        page_numbers = [s.page_number for s in doc.sections]
        assert page_numbers == sorted(page_numbers)

    def test_load_handles_speaker_notes(self, sample_pptx: Path):
        loader = PPTXLoader()
        doc = loader.load(str(sample_pptx))

        slide_two = next(s for s in doc.sections if s.heading == "Slide Two")
        assert "Speaker Notes:" in slide_two.content

    def test_load_detects_images(self, pptx_with_image: Path):
        loader = PPTXLoader()
        doc = loader.load(str(pptx_with_image))

        assert "[Image on slide 1]" in doc.content

    def test_load_handles_empty_presentation(self, empty_pptx: Path):
        loader = PPTXLoader()
        doc = loader.load(str(empty_pptx))

        assert doc.content == "[Empty presentation]"

    def test_load_file_not_found(self):
        loader = PPTXLoader()
        with pytest.raises(FileNotFoundError):
            loader.load("nonexistent.pptx")

    def test_load_invalid_extension(self, tmp_path: Path):
        loader = PPTXLoader()
        txt_file = tmp_path / "bad.txt"
        txt_file.write_text("not pptx")

        with pytest.raises(ValueError, match="not a .pptx file"):
            loader.load(str(txt_file))
